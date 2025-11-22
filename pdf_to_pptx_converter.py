import sys
import os
import platform
import subprocess
from io import BytesIO
from datetime import datetime

# GUI Components
from PyQt5.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QHBoxLayout, QPushButton, 
    QListWidget, QFileDialog, QLabel, QProgressBar, QMessageBox, 
    QAbstractItemView, QListWidgetItem, QFrame, QComboBox
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QSize
from PyQt5.QtGui import QIcon, QDragEnterEvent, QDropEvent

# Processing Libraries
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Cm
from PIL import Image

class ConversionWorker(QThread):
    """
    変換処理をバックグラウンドで行うスレッド
    """
    progress_updated = pyqtSignal(int)
    status_updated = pyqtSignal(str) # ログではなくステータス表示用
    finished_signal = pyqtSignal(str)
    error_signal = pyqtSignal(str)

    def __init__(self, file_list, output_dir, dpi=300):
        super().__init__()
        self.file_list = file_list
        self.output_dir = output_dir
        self.dpi = dpi
        self.is_running = True

    def get_reference_size(self, file_path):
        """基準サイズ取得"""
        try:
            lower_path = file_path.lower()
            if lower_path.endswith('.pdf'):
                images = convert_from_path(file_path, dpi=self.dpi, first_page=1, last_page=1)
                if images:
                    return images[0].size
            else:
                with Image.open(file_path) as img:
                    return img.size
        except Exception:
            pass
        return None

    def run(self):
        total_files = len(self.file_list)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # ログファイル出力は廃止
        
        self.status_updated.emit("初期化中...")

        if not self.file_list:
            return

        # 1. スライドサイズの決定
        first_file = self.file_list[0]
        self.status_updated.emit(f"基準サイズを確認中: {os.path.basename(first_file)}")
        
        prs = Presentation()
        ref_size = self.get_reference_size(first_file)
        
        if ref_size:
            width_px, height_px = ref_size
            # DPIに基づくEMU換算
            prs.slide_width = int(width_px / self.dpi * 914400)
            prs.slide_height = int(height_px / self.dpi * 914400)
        else:
            self.status_updated.emit("サイズ取得失敗。デフォルトサイズを使用します。")

        output_filename = f"Combined_Slides_{timestamp}.pptx"
        save_path = os.path.join(self.output_dir, output_filename)

        processed_count = 0

        # 2. ファイル処理ループ
        for i, file_path in enumerate(self.file_list):
            if not self.is_running:
                break

            file_name = os.path.basename(file_path)
            
            try:
                images_to_process = []

                if file_path.lower().endswith('.pdf'):
                    # PDFの場合
                    self.status_updated.emit(f"[{i+1}/{total_files}] {file_name}: 画像データを読み込み中...")
                    # PDF変換 (重い処理)
                    images = convert_from_path(file_path, dpi=self.dpi, fmt='jpeg')
                    
                    if not images:
                        self.status_updated.emit(f"警告: 画像が取得できませんでした ({file_name})")
                        continue
                    
                    images_to_process = images
                else:
                    # 画像ファイルの場合
                    self.status_updated.emit(f"[{i+1}/{total_files}] {file_name}: 画像を開いています...")
                    try:
                        img = Image.open(file_path)
                        if img.mode in ('RGBA', 'P'):
                            img = img.convert('RGB')
                        images_to_process = [img]
                    except Exception as e:
                        print(f"Error: {e}")
                        continue

                # スライドへの追加処理 (ページごとの進捗を表示)
                total_pages = len(images_to_process)
                for p_idx, img in enumerate(images_to_process):
                    if not self.is_running:
                        break
                    
                    # 詳細なステータス表示: ファイル名 ページ X / Y
                    msg = f"[{i+1}/{total_files}] {file_name}: ページ {p_idx + 1} / {total_pages} を変換中"
                    self.status_updated.emit(msg)

                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)

                    image_stream = BytesIO()
                    jpg_quality = 95 if self.dpi >= 300 else 85
                    img.save(image_stream, format="JPEG", quality=jpg_quality)
                    image_stream.seek(0)

                    slide.shapes.add_picture(
                        image_stream, 
                        0, 0, 
                        width=prs.slide_width, 
                        height=prs.slide_height
                    )
                    image_stream.close()

            except Exception as e:
                self.status_updated.emit(f"エラー発生 ({file_name}): {str(e)}")
            
            processed_count += 1
            progress = int((processed_count / total_files) * 100)
            self.progress_updated.emit(progress)

        # 保存
        if self.is_running:
            self.status_updated.emit("ファイルを保存しています...")
            try:
                prs.save(save_path)
                self.status_updated.emit("保存完了！")
            except Exception as e:
                self.status_updated.emit(f"保存失敗: {str(e)}")

        self.finished_signal.emit(self.output_dir)

    def stop(self):
        self.is_running = False


class FileDropListWidget(QListWidget):
    VALID_EXTENSIONS = ('.pdf', '.png', '.jpg', '.jpeg', '.bmp', '.tiff')

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setAcceptDrops(True)
        self.setSelectionMode(QAbstractItemView.ExtendedSelection)

    def dragEnterEvent(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls():
            event.accept()
        else:
            event.ignore()

    def dragMoveEvent(self, event):
        if event.mimeData().hasUrls():
            event.accept()
        else:
            event.ignore()

    def dropEvent(self, event: QDropEvent):
        if event.mimeData().hasUrls():
            files = [u.toLocalFile() for u in event.mimeData().urls()]
            self.add_files(files)
            event.accept()
        else:
            event.ignore()

    def add_files(self, file_paths):
        for path in file_paths:
            if os.path.isfile(path) and path.lower().endswith(self.VALID_EXTENSIONS):
                items = [self.item(i).text() for i in range(self.count())]
                if path not in items:
                    self.addItem(path)
            elif os.path.isdir(path):
                self.add_folder(path)
    
    def add_folder(self, folder_path):
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                if file.lower().endswith(self.VALID_EXTENSIONS):
                    full_path = os.path.join(root, file)
                    items = [self.item(i).text() for i in range(self.count())]
                    if full_path not in items:
                        self.addItem(full_path)


class PdfToPptApp(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PDF/Image to PPTX Converter")
        self.resize(700, 550)
        self.init_ui()
        self.worker = None

    def init_ui(self):
        main_layout = QVBoxLayout()

        # --- ヘッダー / ボタンエリア ---
        btn_layout = QHBoxLayout()
        
        self.btn_add_file = QPushButton("ファイルを追加")
        self.btn_add_file.clicked.connect(self.open_file_dialog)
        
        self.btn_add_folder = QPushButton("フォルダを追加")
        self.btn_add_folder.clicked.connect(self.open_folder_dialog)
        
        self.btn_clear = QPushButton("リストをクリア / リセット")
        self.btn_clear.clicked.connect(self.clear_list_and_reset) # 関数名を変更

        btn_layout.addWidget(self.btn_add_file)
        btn_layout.addWidget(self.btn_add_folder)
        btn_layout.addStretch()
        btn_layout.addWidget(self.btn_clear)

        main_layout.addLayout(btn_layout)

        # --- メインリストと操作ボタンエリア ---
        list_layout = QHBoxLayout()

        self.list_widget = FileDropListWidget()
        list_layout.addWidget(self.list_widget)

        side_btn_layout = QVBoxLayout()
        
        self.btn_up = QPushButton("↑")
        self.btn_up.setFixedWidth(40)
        self.btn_up.clicked.connect(self.move_item_up)
        
        self.btn_remove = QPushButton("削除")
        self.btn_remove.setFixedWidth(40)
        self.btn_remove.setStyleSheet("color: red;")
        self.btn_remove.clicked.connect(self.remove_selected_item)
        
        self.btn_down = QPushButton("↓")
        self.btn_down.setFixedWidth(40)
        self.btn_down.clicked.connect(self.move_item_down)

        side_btn_layout.addStretch()
        side_btn_layout.addWidget(self.btn_up)
        side_btn_layout.addWidget(self.btn_remove)
        side_btn_layout.addWidget(self.btn_down)
        side_btn_layout.addStretch()

        list_layout.addLayout(side_btn_layout)
        main_layout.addLayout(list_layout)

        lbl_instruction = QLabel("PDFまたは画像(JPG, PNG等)をドラッグ＆ドロップできます。\n最初のファイルのサイズが全PPTXのスライドサイズの基準になります。")
        lbl_instruction.setAlignment(Qt.AlignCenter)
        lbl_instruction.setStyleSheet("color: gray; font-style: italic;")
        main_layout.addWidget(lbl_instruction)
        
        # --- 設定エリア (DPI) ---
        settings_frame = QFrame()
        settings_frame.setFrameShape(QFrame.StyledPanel)
        settings_layout = QHBoxLayout(settings_frame)
        
        lbl_dpi = QLabel("変換画質 (DPI):")
        self.combo_dpi = QComboBox()
        self.combo_dpi.addItem("低画質 / スクリーン用 (96 dpi)", 96)
        self.combo_dpi.addItem("中画質 / 標準 (150 dpi)", 150)
        self.combo_dpi.addItem("高画質 / 印刷用 (300 dpi)", 300)
        self.combo_dpi.addItem("最高画質 (400 dpi)", 400)
        self.combo_dpi.setCurrentIndex(2)

        settings_layout.addWidget(lbl_dpi)
        settings_layout.addWidget(self.combo_dpi)
        settings_layout.addStretch()
        
        main_layout.addWidget(settings_frame)

        # --- 実行エリア ---
        action_layout = QVBoxLayout()
        
        # ステータス表示エリアを強調
        self.status_label = QLabel("待機中...")
        self.status_label.setStyleSheet("font-weight: bold; color: #333; margin-bottom: 5px;")
        action_layout.addWidget(self.status_label)

        self.progress_bar = QProgressBar()
        self.progress_bar.setValue(0)
        action_layout.addWidget(self.progress_bar)

        self.btn_convert = QPushButton("変換開始")
        self.btn_convert.setMinimumHeight(50)
        self.btn_convert.setStyleSheet("font-weight: bold; font-size: 14px;")
        self.btn_convert.clicked.connect(self.start_conversion)
        action_layout.addWidget(self.btn_convert)

        main_layout.addLayout(action_layout)

        self.setLayout(main_layout)

    # --- ロジック ---

    def open_file_dialog(self):
        filters = "Supported Files (*.pdf *.png *.jpg *.jpeg *.bmp *.tiff);;PDF Files (*.pdf);;Images (*.png *.jpg *.jpeg *.bmp *.tiff)"
        files, _ = QFileDialog.getOpenFileNames(self, "ファイルを選択", "", filters)
        if files:
            self.list_widget.add_files(files)

    def open_folder_dialog(self):
        folder = QFileDialog.getExistingDirectory(self, "フォルダを選択")
        if folder:
            self.list_widget.add_folder(folder)

    def clear_list_and_reset(self):
        """リストをクリアし、すべての状態を初期化する"""
        # リストのクリア
        self.list_widget.clear()
        
        # 状態のリセット
        self.progress_bar.setValue(0)
        self.status_label.setText("待機中... (リセット完了)")
        
        # ワーカーのクリーンアップ（もし存在すれば）
        if self.worker and self.worker.isRunning():
            self.worker.stop()
            self.worker.wait()
        self.worker = None
        
        # UIのロック解除
        self.toggle_ui(True)

    def remove_selected_item(self):
        selected_items = self.list_widget.selectedItems()
        if not selected_items:
            return
        for item in selected_items:
            self.list_widget.takeItem(self.list_widget.row(item))

    def move_item_up(self):
        curr_row = self.list_widget.currentRow()
        if curr_row > 0:
            item = self.list_widget.takeItem(curr_row)
            self.list_widget.insertItem(curr_row - 1, item)
            self.list_widget.setCurrentRow(curr_row - 1)

    def move_item_down(self):
        curr_row = self.list_widget.currentRow()
        if curr_row < self.list_widget.count() - 1 and curr_row >= 0:
            item = self.list_widget.takeItem(curr_row)
            self.list_widget.insertItem(curr_row + 1, item)
            self.list_widget.setCurrentRow(curr_row + 1)

    def start_conversion(self):
        count = self.list_widget.count()
        if count == 0:
            QMessageBox.warning(self, "警告", "ファイルが選択されていません。")
            return

        output_dir = QFileDialog.getExistingDirectory(self, "保存先フォルダを選択")
        if not output_dir:
            return
            
        selected_dpi = self.combo_dpi.currentData()

        self.toggle_ui(False)
        self.progress_bar.setValue(0)
        self.status_label.setText("準備中...")

        file_list = [self.list_widget.item(i).text() for i in range(count)]

        self.worker = ConversionWorker(file_list, output_dir, dpi=selected_dpi)
        self.worker.progress_updated.connect(self.update_progress)
        self.worker.status_updated.connect(self.update_status) # ログではなくステータス更新に接続
        self.worker.finished_signal.connect(self.conversion_finished)
        self.worker.start()

    def update_progress(self, value):
        self.progress_bar.setValue(value)

    def update_status(self, message):
        """ステータスラベルを更新"""
        self.status_label.setText(message)

    def conversion_finished(self, output_dir):
        self.toggle_ui(True)
        self.progress_bar.setValue(100)
        self.status_label.setText("変換完了！")
        
        QMessageBox.information(self, "完了", "すべてのファイルの処理が完了しました。")
        
        self.open_folder(output_dir)

    def open_folder(self, path):
        try:
            if platform.system() == "Windows":
                os.startfile(path)
            elif platform.system() == "Darwin":
                subprocess.run(["open", path])
            else:
                subprocess.run(["xdg-open", path])
        except Exception as e:
            self.status_label.setText(f"フォルダを開けませんでした: {str(e)}")

    def toggle_ui(self, enabled):
        self.btn_add_file.setEnabled(enabled)
        self.btn_add_folder.setEnabled(enabled)
        self.btn_clear.setEnabled(enabled)
        self.btn_remove.setEnabled(enabled)
        self.btn_up.setEnabled(enabled)
        self.btn_down.setEnabled(enabled)
        self.list_widget.setEnabled(enabled)
        self.combo_dpi.setEnabled(enabled)
        
        if enabled:
            self.btn_convert.setText("変換開始")
            # シグナル接続の重複を防ぐため一度切断
            try:
                self.btn_convert.clicked.disconnect()
            except TypeError:
                pass # 未接続なら無視
            self.btn_convert.clicked.connect(self.start_conversion)
            
            # ★修正箇所: ボタンを有効化する
            self.btn_convert.setEnabled(True) 
        else:
            self.btn_convert.setText("処理中...")
            self.btn_convert.setEnabled(False) 

if __name__ == "__main__":
    app = QApplication(sys.argv)
    
    if hasattr(Qt, 'AA_EnableHighDpiScaling'):
        QApplication.setAttribute(Qt.AA_EnableHighDpiScaling, True)
    if hasattr(Qt, 'AA_UseHighDpiPixmaps'):
        QApplication.setAttribute(Qt.AA_UseHighDpiPixmaps, True)

    window = PdfToPptApp()
    window.show()
    sys.exit(app.exec_())